#!/usr/bin/env python3
"""
build-pptx-template.py

Generates the cytario branded PowerPoint template from design system tokens.
Produces assets/templates/cytario-template.pptx with six slide layouts:

  1. Title Slide          - Full-bleed purple background, logo, title, subtitle
  2. Section Divider      - Teal accent bar, section title
  3. Content Slide        - Title + body text on white
  4. Two-Column Layout    - Title + two equal columns
  5. Image + Text Layout  - Left image area + right text
  6. Closing Slide        - Purple background, logo, thank-you message

Brand elements:
  - Colors from tokens/base.json (purple-700 #5c2483, teal-500 #35b7b8)
  - Typography: Montserrat (must be installed on the system)
  - Logo: cytario-logo-mono-white.png for dark backgrounds,
          cytario-logo-purple.png for light backgrounds
"""

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Brand constants (from tokens/base.json and tokens/semantic.json)
# ---------------------------------------------------------------------------
PURPLE_700 = RGBColor(0x5C, 0x24, 0x83)  # Brand primary
PURPLE_800 = RGBColor(0x4A, 0x1D, 0x6A)
PURPLE_900 = RGBColor(0x3A, 0x17, 0x54)
PURPLE_100 = RGBColor(0xEA, 0xD9, 0xF5)
PURPLE_50 = RGBColor(0xF5, 0xF0, 0xFA)
TEAL_500 = RGBColor(0x35, 0xB7, 0xB8)    # Brand accent
TEAL_700 = RGBColor(0x21, 0x7D, 0x7E)    # Accessible teal for text
TEAL_100 = RGBColor(0xD0, 0xF0, 0xF0)
TEAL_50 = RGBColor(0xED, 0xF9, 0xF9)
NEUTRAL_0 = RGBColor(0xFF, 0xFF, 0xFF)    # White
NEUTRAL_50 = RGBColor(0xF9, 0xFA, 0xFB)
NEUTRAL_100 = RGBColor(0xF3, 0xF4, 0xF6)
NEUTRAL_200 = RGBColor(0xE5, 0xE7, 0xEB)
NEUTRAL_400 = RGBColor(0x9C, 0xA3, 0xAF)
NEUTRAL_500 = RGBColor(0x6B, 0x72, 0x80)
NEUTRAL_600 = RGBColor(0x4B, 0x55, 0x63)
NEUTRAL_700 = RGBColor(0x37, 0x41, 0x51)
NEUTRAL_800 = RGBColor(0x1F, 0x29, 0x37)
NEUTRAL_900 = RGBColor(0x11, 0x18, 0x27)

# Typography
FONT_FAMILY = "Montserrat"

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Paths
PROJECT_ROOT = Path(__file__).resolve().parent.parent
LOGO_DIR = PROJECT_ROOT / "assets" / "logos" / "png"
TEMPLATE_DIR = PROJECT_ROOT / "assets" / "templates"
OUTPUT_PATH = TEMPLATE_DIR / "cytario-template.pptx"

LOGO_WHITE = LOGO_DIR / "cytario-logo-mono-white.png"
LOGO_DEFAULT = LOGO_DIR / "cytario-logo-purple.png"
LOGO_ON_PURPLE = LOGO_DIR / "cytario-logo-on-purple.png"
LOGO_FULL = LOGO_DIR / "cytario-logo-purple-tagline.png"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def set_slide_bg(slide, color):
    """Set a solid fill background color on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rectangle(slide, left, top, width, height, color, transparency=0):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # No border
    if transparency > 0:
        # Set transparency via alpha (0 = opaque, 100000 = fully transparent in OOXML units)
        shape.fill.fore_color.brightness = 0
    return shape


def set_text(textframe, text, font_size=Pt(16), color=NEUTRAL_900,
             bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT_FAMILY,
             spacing_before=Pt(0), spacing_after=Pt(0), line_spacing=1.25):
    """Set text in a textframe with brand typography."""
    textframe.clear()
    textframe.word_wrap = True
    p = textframe.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = spacing_before
    p.space_after = spacing_after
    if line_spacing:
        p.line_spacing = line_spacing
    return p


def add_textbox(slide, left, top, width, height, text, font_size=Pt(16),
                color=NEUTRAL_900, bold=False, alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font_name=FONT_FAMILY):
    """Add a text box with brand typography."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    txbox.text_frame.word_wrap = True
    tf = txbox.text_frame
    tf.auto_size = None
    # Set vertical alignment via anchor
    # (textframe doesn't have anchor directly; we set it on the shape)
    set_text(tf, text, font_size=font_size, color=color, bold=bold,
             alignment=alignment, font_name=font_name)
    return txbox


def add_footer_bar(slide, text="cytario  |  Connecting Pathology", color=NEUTRAL_400):
    """Add a subtle footer bar at the bottom of a content slide."""
    # Thin teal line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.75), SLIDE_HEIGHT - Inches(0.65),
        SLIDE_WIDTH - Inches(1.5), Pt(1.5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL_500
    line.line.fill.background()

    # Footer text
    add_textbox(
        slide,
        Inches(0.75), SLIDE_HEIGHT - Inches(0.6),
        Inches(6), Inches(0.35),
        text,
        font_size=Pt(9), color=NEUTRAL_400,
        alignment=PP_ALIGN.LEFT
    )


def add_slide_number_placeholder(slide, color=NEUTRAL_400):
    """Add a page number text in the bottom-right corner."""
    add_textbox(
        slide,
        SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.6),
        Inches(0.75), Inches(0.35),
        "#",
        font_size=Pt(9), color=color,
        alignment=PP_ALIGN.RIGHT
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_title_slide(prs):
    """
    Layout 1: Title Slide
    Full purple background, white logo top-left, teal accent line,
    large title + subtitle in white, decorative teal corner element.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, PURPLE_700)

    # Decorative teal accent bar at top
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(6), TEAL_500)

    # Decorative teal rectangle in bottom-right corner (abstract geometric element)
    add_rectangle(
        slide,
        SLIDE_WIDTH - Inches(4.5), SLIDE_HEIGHT - Inches(2.2),
        Inches(4.5), Inches(2.2),
        PURPLE_800
    )

    # Teal accent line (vertical, left side)
    add_rectangle(
        slide,
        Inches(0.75), Inches(2.5),
        Pt(4), Inches(2.8),
        TEAL_500
    )

    # Logo (white variant, top-left)
    if LOGO_ON_PURPLE.exists():
        slide.shapes.add_picture(
            str(LOGO_ON_PURPLE),
            Inches(0.75), Inches(0.6),
            height=Inches(0.7)
        )

    # Title text
    add_textbox(
        slide,
        Inches(1.2), Inches(2.5),
        Inches(8.5), Inches(1.5),
        "Presentation Title",
        font_size=Pt(40), color=NEUTRAL_0, bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Subtitle text
    add_textbox(
        slide,
        Inches(1.2), Inches(4.1),
        Inches(8), Inches(0.8),
        "Subtitle or description line",
        font_size=Pt(20), color=TEAL_100, bold=False,
        alignment=PP_ALIGN.LEFT
    )

    # Date / presenter info
    add_textbox(
        slide,
        Inches(1.2), Inches(5.2),
        Inches(6), Inches(0.5),
        "Presenter Name  |  Date  |  Event",
        font_size=Pt(14), color=PURPLE_100, bold=False,
        alignment=PP_ALIGN.LEFT
    )

    # Tagline bottom-right
    add_textbox(
        slide,
        SLIDE_WIDTH - Inches(4), SLIDE_HEIGHT - Inches(0.9),
        Inches(3.5), Inches(0.5),
        "Connecting Pathology",
        font_size=Pt(14), color=TEAL_500, bold=True,
        alignment=PP_ALIGN.RIGHT
    )


def build_section_divider(prs):
    """
    Layout 2: Section Divider
    White background with a bold purple left band and teal accent.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEUTRAL_0)

    # Purple left band
    add_rectangle(slide, Inches(0), Inches(0), Inches(0.6), SLIDE_HEIGHT, PURPLE_700)

    # Teal accent strip on top of the purple band
    add_rectangle(slide, Inches(0), Inches(0), Inches(0.6), Pt(6), TEAL_500)

    # Light purple background area for title
    add_rectangle(
        slide,
        Inches(0.6), Inches(2.0),
        SLIDE_WIDTH - Inches(0.6), Inches(3.5),
        PURPLE_50
    )

    # Teal accent line above title
    add_rectangle(
        slide,
        Inches(1.5), Inches(2.8),
        Inches(2.5), Pt(4),
        TEAL_500
    )

    # Section number
    add_textbox(
        slide,
        Inches(1.5), Inches(2.2),
        Inches(2), Inches(0.5),
        "01",
        font_size=Pt(16), color=TEAL_700, bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Section title
    add_textbox(
        slide,
        Inches(1.5), Inches(3.1),
        Inches(9), Inches(1.2),
        "Section Title",
        font_size=Pt(36), color=PURPLE_700, bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Section subtitle
    add_textbox(
        slide,
        Inches(1.5), Inches(4.3),
        Inches(8), Inches(0.6),
        "Brief description of what this section covers",
        font_size=Pt(16), color=NEUTRAL_600, bold=False,
        alignment=PP_ALIGN.LEFT
    )

    # Footer
    add_footer_bar(slide)
    add_slide_number_placeholder(slide)


def build_content_slide(prs):
    """
    Layout 3: Content Slide
    White background, title at top, body text area below.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEUTRAL_0)

    # Teal top accent line
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), TEAL_500)

    # Logo (small, top-right)
    if LOGO_DEFAULT.exists():
        slide.shapes.add_picture(
            str(LOGO_DEFAULT),
            SLIDE_WIDTH - Inches(2.2), Inches(0.35),
            height=Inches(0.35)
        )

    # Slide title
    add_textbox(
        slide,
        Inches(0.75), Inches(0.5),
        Inches(9), Inches(0.8),
        "Slide Title",
        font_size=Pt(28), color=PURPLE_700, bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Thin separator line under title
    add_rectangle(
        slide,
        Inches(0.75), Inches(1.35),
        Inches(2), Pt(2.5),
        TEAL_500
    )

    # Body text area
    body_box = add_textbox(
        slide,
        Inches(0.75), Inches(1.7),
        Inches(11.5), Inches(4.8),
        "Body text goes here. Use bullet points or paragraphs to present your content clearly and concisely.",
        font_size=Pt(16), color=NEUTRAL_700, bold=False,
        alignment=PP_ALIGN.LEFT
    )

    # Footer
    add_footer_bar(slide)
    add_slide_number_placeholder(slide)


def build_two_column_slide(prs):
    """
    Layout 4: Two-Column Layout
    Title at top, two equal columns below.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEUTRAL_0)

    # Teal top accent line
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), TEAL_500)

    # Logo (small, top-right)
    if LOGO_DEFAULT.exists():
        slide.shapes.add_picture(
            str(LOGO_DEFAULT),
            SLIDE_WIDTH - Inches(2.2), Inches(0.35),
            height=Inches(0.35)
        )

    # Slide title
    add_textbox(
        slide,
        Inches(0.75), Inches(0.5),
        Inches(9), Inches(0.8),
        "Two-Column Title",
        font_size=Pt(28), color=PURPLE_700, bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Thin separator line under title
    add_rectangle(
        slide,
        Inches(0.75), Inches(1.35),
        Inches(2), Pt(2.5),
        TEAL_500
    )

    col_width = Inches(5.5)
    col_top = Inches(1.8)
    col_height = Inches(4.5)
    left_margin = Inches(0.75)
    gap = Inches(0.5)

    # Left column header
    add_textbox(
        slide,
        left_margin, col_top,
        col_width, Inches(0.5),
        "Column One",
        font_size=Pt(18), color=PURPLE_700, bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Left column teal accent
    add_rectangle(
        slide,
        left_margin, col_top + Inches(0.55),
        Inches(1.5), Pt(2.5),
        TEAL_500
    )

    # Left column body
    add_textbox(
        slide,
        left_margin, col_top + Inches(0.8),
        col_width, col_height - Inches(0.8),
        "Content for the left column. Use this space for key points, data, or explanations.",
        font_size=Pt(14), color=NEUTRAL_700, bold=False,
        alignment=PP_ALIGN.LEFT
    )

    # Right column header
    right_left = left_margin + col_width + gap
    add_textbox(
        slide,
        right_left, col_top,
        col_width, Inches(0.5),
        "Column Two",
        font_size=Pt(18), color=PURPLE_700, bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Right column teal accent
    add_rectangle(
        slide,
        right_left, col_top + Inches(0.55),
        Inches(1.5), Pt(2.5),
        TEAL_500
    )

    # Right column body
    add_textbox(
        slide,
        right_left, col_top + Inches(0.8),
        col_width, col_height - Inches(0.8),
        "Content for the right column. Use this space for supporting information, comparisons, or details.",
        font_size=Pt(14), color=NEUTRAL_700, bold=False,
        alignment=PP_ALIGN.LEFT
    )

    # Vertical divider line between columns
    add_rectangle(
        slide,
        left_margin + col_width + Inches(0.22), col_top + Inches(0.2),
        Pt(1), col_height - Inches(0.4),
        NEUTRAL_200
    )

    # Footer
    add_footer_bar(slide)
    add_slide_number_placeholder(slide)


def build_image_text_slide(prs):
    """
    Layout 5: Image + Text
    Left half for an image placeholder, right half for text.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NEUTRAL_0)

    # Teal top accent line
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(4), TEAL_500)

    # Logo (small, top-right)
    if LOGO_DEFAULT.exists():
        slide.shapes.add_picture(
            str(LOGO_DEFAULT),
            SLIDE_WIDTH - Inches(2.2), Inches(0.35),
            height=Inches(0.35)
        )

    # Left image placeholder area (light gray background)
    img_area = add_rectangle(
        slide,
        Inches(0.75), Inches(0.75),
        Inches(5.8), Inches(5.6),
        NEUTRAL_100
    )

    # Placeholder text centered in image area
    add_textbox(
        slide,
        Inches(1.5), Inches(3.2),
        Inches(4.3), Inches(0.8),
        "Insert Image Here",
        font_size=Pt(18), color=NEUTRAL_400, bold=False,
        alignment=PP_ALIGN.CENTER
    )

    # Teal border accent on left of image area
    add_rectangle(
        slide,
        Inches(0.75), Inches(0.75),
        Pt(4), Inches(5.6),
        TEAL_500
    )

    # Right side: title
    right_left = Inches(7.0)
    right_width = Inches(5.5)

    add_textbox(
        slide,
        right_left, Inches(1.2),
        right_width, Inches(0.8),
        "Image Title",
        font_size=Pt(28), color=PURPLE_700, bold=True,
        alignment=PP_ALIGN.LEFT
    )

    # Teal accent under title
    add_rectangle(
        slide,
        right_left, Inches(2.0),
        Inches(2), Pt(2.5),
        TEAL_500
    )

    # Right side: body text
    add_textbox(
        slide,
        right_left, Inches(2.3),
        right_width, Inches(3.8),
        "Describe the image or figure shown on the left. Use this space for captions, data interpretation, or key takeaways.\n\nBullet points work well here:\n- First observation\n- Second observation\n- Third observation",
        font_size=Pt(14), color=NEUTRAL_700, bold=False,
        alignment=PP_ALIGN.LEFT
    )

    # Footer
    add_footer_bar(slide)
    add_slide_number_placeholder(slide)


def build_closing_slide(prs):
    """
    Layout 6: Closing / Thank You Slide
    Purple background with centered logo and message.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PURPLE_700)

    # Teal accent bar at top
    add_rectangle(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(6), TEAL_500)

    # Decorative darker purple area at bottom
    add_rectangle(
        slide,
        Inches(0), SLIDE_HEIGHT - Inches(2),
        SLIDE_WIDTH, Inches(2),
        PURPLE_800
    )

    # Logo (centered, white variant)
    if LOGO_ON_PURPLE.exists():
        logo_width = Inches(3.0)
        slide.shapes.add_picture(
            str(LOGO_ON_PURPLE),
            (SLIDE_WIDTH - logo_width) // 2, Inches(1.8),
            width=logo_width
        )

    # Thank you text
    add_textbox(
        slide,
        Inches(2), Inches(3.4),
        SLIDE_WIDTH - Inches(4), Inches(1.2),
        "Thank You",
        font_size=Pt(44), color=NEUTRAL_0, bold=True,
        alignment=PP_ALIGN.CENTER
    )

    # Contact / next steps
    add_textbox(
        slide,
        Inches(2), Inches(4.6),
        SLIDE_WIDTH - Inches(4), Inches(0.6),
        "Questions and discussion",
        font_size=Pt(18), color=TEAL_100, bold=False,
        alignment=PP_ALIGN.CENTER
    )

    # Contact info line
    add_textbox(
        slide,
        Inches(2), SLIDE_HEIGHT - Inches(1.5),
        SLIDE_WIDTH - Inches(4), Inches(0.8),
        "www.cytario.com  |  info@cytario.com",
        font_size=Pt(14), color=PURPLE_100, bold=False,
        alignment=PP_ALIGN.CENTER
    )

    # Tagline
    add_textbox(
        slide,
        Inches(2), SLIDE_HEIGHT - Inches(0.7),
        SLIDE_WIDTH - Inches(4), Inches(0.4),
        "Connecting Pathology",
        font_size=Pt(12), color=TEAL_500, bold=True,
        alignment=PP_ALIGN.CENTER
    )


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Building cytario PPTX template...")
    print(f"  Slide dimensions: {SLIDE_WIDTH} x {SLIDE_HEIGHT} (16:9)")
    print(f"  Logo dir: {LOGO_DIR}")
    print(f"  Output: {OUTPUT_PATH}")

    build_title_slide(prs)
    print("  [1/6] Title slide")

    build_section_divider(prs)
    print("  [2/6] Section divider")

    build_content_slide(prs)
    print("  [3/6] Content slide")

    build_two_column_slide(prs)
    print("  [4/6] Two-column layout")

    build_image_text_slide(prs)
    print("  [5/6] Image + text layout")

    build_closing_slide(prs)
    print("  [6/6] Closing slide")

    prs.save(str(OUTPUT_PATH))
    print(f"\nTemplate saved to: {OUTPUT_PATH}")
    print(f"File size: {OUTPUT_PATH.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
